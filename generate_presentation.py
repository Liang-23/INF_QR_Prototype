"""
generate_presentation.py

One-off generator for the course-presentation deck about the
INF QR Environmental Monitor prototype.

Run:
    python generate_presentation.py

Output:
    presentation.pptx  (next to this file)

Style: "Teal Trust" palette, 16:9, Cambria headers + Calibri body.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- palette ----------
TEAL        = RGBColor(0x02, 0x80, 0x90)   # primary
SEAFOAM     = RGBColor(0x00, 0xA8, 0x96)   # secondary
MINT        = RGBColor(0x02, 0xC3, 0x9A)   # accent
DARK_BG     = RGBColor(0x07, 0x2A, 0x32)   # dark teal for title/closing
TEXT_DARK   = RGBColor(0x1A, 0x2B, 0x33)
TEXT_MUTED  = RGBColor(0x6B, 0x7B, 0x82)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BG     = RGBColor(0xEC, 0xF3, 0xF4)   # very pale teal for card surfaces
CARD_LINE   = RGBColor(0xCF, 0xDE, 0xE1)
CODE_BG     = RGBColor(0x0F, 0x1F, 0x26)
CODE_TXT    = RGBColor(0xC8, 0xE6, 0xC9)

HEADER_FONT = "Cambria"
BODY_FONT   = "Calibri"
MONO_FONT   = "Consolas"

# ---------- helpers ----------

def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def set_text(text_frame, text, *, font=BODY_FONT, size=14, bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False):
    text_frame.margin_left = Inches(0)
    text_frame.margin_right = Inches(0)
    text_frame.margin_top = Inches(0)
    text_frame.margin_bottom = Inches(0)
    text_frame.word_wrap = True
    text_frame.vertical_anchor = anchor
    p = text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def add_textbox(slide, x, y, w, h, text, **kwargs):
    tb = slide.shapes.add_textbox(x, y, w, h)
    set_text(tb.text_frame, text, **kwargs)
    return tb


def add_paragraphs(slide, x, y, w, h, paragraphs, *, font=BODY_FONT,
                   size=14, color=TEXT_DARK, align=PP_ALIGN.LEFT,
                   line_spacing=1.2):
    """paragraphs is a list of (text, bold) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    for i, item in enumerate(paragraphs):
        text, bold = item if isinstance(item, tuple) else (item, False)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    fill_solid(bg, color)
    no_line(bg)
    return bg


def add_card(slide, x, y, w, h, *, fill=SOFT_BG, line=CARD_LINE,
             radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(radius_shape, x, y, w, h)
    fill_solid(s, fill)
    s.line.color.rgb = line
    s.line.width = Pt(0.75)
    # Soften the round corner
    try:
        s.adjustments[0] = 0.08
    except Exception:
        pass
    return s


def add_circle(slide, x, y, d, *, color=MINT):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    fill_solid(s, color)
    no_line(s)
    return s


def add_numbered_circle(slide, cx_in, cy_in, d_in, number, *,
                        color=MINT, text_color=WHITE, size=20):
    d = Inches(d_in)
    x = Inches(cx_in) - d // 2
    y = Inches(cy_in) - d // 2
    s = add_circle(slide, x, y, d, color=color)
    tf = s.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.name = HEADER_FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    return s


def add_title(slide, text, *, color=TEXT_DARK, top=0.55, size=34):
    add_textbox(slide, Inches(0.6), Inches(top), Inches(12.1), Inches(0.7),
                text, font=HEADER_FONT, size=size, bold=True, color=color)


def add_eyebrow(slide, text, *, color=TEAL, top=0.35):
    add_textbox(slide, Inches(0.6), Inches(top), Inches(12.1), Inches(0.3),
                text.upper(), font=BODY_FONT, size=11, bold=True,
                color=color)


def add_footer(slide, page_number, total, *, color=TEXT_MUTED):
    add_textbox(slide, Inches(0.6), Inches(7.05), Inches(6),
                Inches(0.3),
                "INF QR Environmental Monitor",
                font=BODY_FONT, size=10, color=color)
    add_textbox(slide, Inches(11.7), Inches(7.05), Inches(1.0),
                Inches(0.3), f"{page_number} / {total}",
                font=BODY_FONT, size=10, color=color,
                align=PP_ALIGN.RIGHT)


def code_block(slide, x_in, y_in, w_in, h_in, lines):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x_in), Inches(y_in),
                                 Inches(w_in), Inches(h_in))
    fill_solid(box, CODE_BG)
    no_line(box)
    try:
        box.adjustments[0] = 0.06
    except Exception:
        pass
    tb = slide.shapes.add_textbox(Inches(x_in + 0.2), Inches(y_in + 0.15),
                                  Inches(w_in - 0.4), Inches(h_in - 0.3))
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = line
        run.font.name = MONO_FONT
        run.font.size = Pt(11)
        run.font.color.rgb = CODE_TXT


# ---------- presentation ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

TOTAL = 11  # update when slide count changes


# ---------- Slide 1: title ----------
s = prs.slides.add_slide(blank_layout)
add_bg(s, DARK_BG)
# Accent corner marks (mint circle + seafoam ring)
add_circle(s, Inches(11.3), Inches(0.7), Inches(0.45), color=MINT)
ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.1), Inches(0.55),
                          Inches(0.75), Inches(0.75))
ring.fill.background()
ring.line.color.rgb = SEAFOAM
ring.line.width = Pt(2)

add_textbox(s, Inches(0.8), Inches(1.4), Inches(12), Inches(0.5),
            "COURSE PRESENTATION  ·  IoT PROTOTYPE",
            font=BODY_FONT, size=14, bold=True, color=MINT)
add_textbox(s, Inches(0.8), Inches(2.1), Inches(12), Inches(1.6),
            "INF QR Environmental Monitor",
            font=HEADER_FONT, size=54, bold=True, color=WHITE)
add_textbox(s, Inches(0.8), Inches(3.6), Inches(11.5), Inches(1.2),
            "Live temperature, humidity, and light data delivered to a phone "
            "through a dynamic QR code — no app install required.",
            font=BODY_FONT, size=20, color=RGBColor(0xCA, 0xDC, 0xDC),
            italic=True)

# Divider
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.3),
                         Inches(0.6), Inches(0.06))
fill_solid(div, MINT)
no_line(div)

add_textbox(s, Inches(0.8), Inches(5.5), Inches(12), Inches(0.4),
            "Arduino UNO R3   ·   Flask Server   ·   Magic QR Refresh   ·   I²C LED Matrix",
            font=BODY_FONT, size=14, color=RGBColor(0x9C, 0xC2, 0xC6))

add_textbox(s, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
            "Prototype demo · 2026",
            font=BODY_FONT, size=12, color=SEAFOAM)


# ---------- Slide 2: problem & motivation ----------
s = prs.slides.add_slide(blank_layout)
add_bg(s, WHITE)
add_eyebrow(s, "Why this project")
add_title(s, "Environmental data on any phone, in one scan")

# Big stat callout (left)
stat_card = add_card(s, Inches(0.6), Inches(1.7), Inches(4.5), Inches(4.5),
                     fill=SOFT_BG)
add_textbox(s, Inches(0.85), Inches(1.95), Inches(4.0), Inches(0.4),
            "ZERO", font=BODY_FONT, size=14, bold=True, color=TEAL)
add_textbox(s, Inches(0.85), Inches(2.4), Inches(4.0), Inches(2.4),
            "0", font=HEADER_FONT, size=180, bold=True, color=TEAL,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(0.85), Inches(5.1), Inches(4.0), Inches(0.5),
            "apps installed", font=HEADER_FONT, size=22, bold=True,
            color=TEXT_DARK)
add_textbox(s, Inches(0.85), Inches(5.55), Inches(4.0), Inches(0.5),
            "on the viewer's phone", font=BODY_FONT, size=14,
            color=TEXT_MUTED)

# Three pillars on the right
pillars = [
    ("Universal access",
     "Any phone with a browser and a camera can view live data."),
    ("Live, not stored",
     "Sensors push fresh readings every three seconds over USB serial."),
    ("Local, low-cost hardware",
     "Built with an Arduino UNO R3 — no Wi-Fi MCU, no cloud account."),
]
top = 1.8
for title, body in pillars:
    add_circle(s, Inches(5.7), Inches(top + 0.15), Inches(0.35),
               color=MINT)
    add_textbox(s, Inches(6.25), Inches(top), Inches(6.7), Inches(0.4),
                title, font=HEADER_FONT, size=20, bold=True,
                color=TEAL)
    add_textbox(s, Inches(6.25), Inches(top + 0.5), Inches(6.7),
                Inches(0.9), body, font=BODY_FONT, size=15,
                color=TEXT_DARK)
    top += 1.45

add_footer(s, 2, TOTAL)


# ---------- Slide 3: system architecture ----------
s = prs.slides.add_slide(blank_layout)
add_bg(s, WHITE)
add_eyebrow(s, "System architecture")
add_title(s, "Sensors → Arduino → Laptop → Phone")

stages = [
    ("Sensors",     "DHT11 + TEMT6000",       MINT),
    ("Arduino",     "UNO R3 · serial JSON",   SEAFOAM),
    ("Laptop",      "Flask gateway (HTTP)",   TEAL),
    ("QR + Wi-Fi",  "LAN /qr endpoint",       SEAFOAM),
    ("Phone",       "Mobile dashboard",       MINT),
]

# Stage boxes — pipeline
n = len(stages)
box_w = 2.15
gap = 0.25
total_w = n * box_w + (n - 1) * gap
left0 = (13.333 - total_w) / 2
top_stage = 2.4
box_h = 2.0

for i, (title, sub, color) in enumerate(stages):
    x = left0 + i * (box_w + gap)
    # Card
    card = add_card(s, Inches(x), Inches(top_stage),
                    Inches(box_w), Inches(box_h),
                    fill=WHITE, line=color)
    card.line.width = Pt(2)
    # Color tab on top
    tab = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(x + 0.18), Inches(top_stage + 0.22),
                             Inches(0.35), Inches(0.08))
    fill_solid(tab, color); no_line(tab)
    # Step number
    add_textbox(s, Inches(x + 0.18), Inches(top_stage + 0.35),
                Inches(1.0), Inches(0.3),
                f"STEP {i+1}", font=BODY_FONT, size=10, bold=True,
                color=color)
    add_textbox(s, Inches(x + 0.18), Inches(top_stage + 0.7),
                Inches(box_w - 0.36), Inches(0.6),
                title, font=HEADER_FONT, size=20, bold=True,
                color=TEXT_DARK)
    add_textbox(s, Inches(x + 0.18), Inches(top_stage + 1.2),
                Inches(box_w - 0.36), Inches(0.7),
                sub, font=BODY_FONT, size=12, color=TEXT_MUTED)
    # Arrow chevron between boxes
    if i < n - 1:
        arrow_x = x + box_w + 0.02
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(arrow_x),
                                   Inches(top_stage + box_h/2 - 0.18),
                                   Inches(gap - 0.04), Inches(0.36))
        fill_solid(arrow, TEAL); no_line(arrow)

# Caption
add_textbox(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.4),
            "The laptop acts as a temporary gateway because the UNO R3 has no built-in Wi-Fi.",
            font=BODY_FONT, size=14, italic=True, color=TEXT_MUTED,
            align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.4),
            "Data updates every 3 seconds end-to-end.",
            font=BODY_FONT, size=14, italic=True, color=TEXT_MUTED,
            align=PP_ALIGN.CENTER)
add_footer(s, 3, TOTAL)


# ---------- Slide 4: hardware ----------
s = prs.slides.add_slide(blank_layout)
add_bg(s, WHITE)
add_eyebrow(s, "Hardware")
add_title(s, "Four parts, one breadboard")

# Components table on left
components = [
    ("Arduino UNO R3", "Main controller", "USB to laptop"),
    ("DHT11",          "Temperature + humidity", "Data → D2"),
    ("TEMT6000",       "Ambient light sensor",   "Signal → A0"),
    ("8×8 LED Matrix", "I²C local feedback",     "SDA → A4, SCL → A5"),
]

tbl_top = 1.8
row_h = 0.95
row_w = 7.4

# Header
header = add_card(s, Inches(0.6), Inches(tbl_top), Inches(row_w),
                  Inches(0.55), fill=TEAL, line=TEAL)
header.line.width = Pt(0)
for col_x, col_w, text in [
    (0.8,  2.2, "Component"),
    (3.0,  2.6, "Role"),
    (5.6,  2.4, "Connection"),
]:
    add_textbox(s, Inches(col_x), Inches(tbl_top + 0.13),
                Inches(col_w), Inches(0.35), text,
                font=BODY_FONT, size=12, bold=True, color=WHITE)

for i, (comp, role, conn) in enumerate(components):
    y = tbl_top + 0.55 + i * row_h
    band_fill = SOFT_BG if i % 2 == 0 else WHITE
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y),
                              Inches(row_w), Inches(row_h))
    fill_solid(band, band_fill)
    band.line.color.rgb = CARD_LINE
    band.line.width = Pt(0.5)
    add_textbox(s, Inches(0.8), Inches(y + 0.2), Inches(2.2), Inches(0.6),
                comp, font=HEADER_FONT, size=16, bold=True,
                color=TEXT_DARK)
    add_textbox(s, Inches(3.0), Inches(y + 0.2), Inches(2.6), Inches(0.6),
                role, font=BODY_FONT, size=13, color=TEXT_DARK)
    add_textbox(s, Inches(5.6), Inches(y + 0.2), Inches(2.4), Inches(0.6),
                conn, font=MONO_FONT, size=12, color=TEAL)

# Wiring callout card on right
add_card(s, Inches(8.4), Inches(1.8), Inches(4.4), Inches(4.5),
         fill=SOFT_BG)
add_textbox(s, Inches(8.7), Inches(2.0), Inches(4.0), Inches(0.4),
            "POWER", font=BODY_FONT, size=11, bold=True, color=TEAL)
add_textbox(s, Inches(8.7), Inches(2.4), Inches(4.0), Inches(0.5),
            "All modules: 5 V + GND",
            font=HEADER_FONT, size=18, bold=True, color=TEXT_DARK)

add_textbox(s, Inches(8.7), Inches(3.2), Inches(4.0), Inches(0.4),
            "I²C ADDRESS", font=BODY_FONT, size=11, bold=True,
            color=TEAL)
add_textbox(s, Inches(8.7), Inches(3.6), Inches(4.0), Inches(0.5),
            "0x70  (HT16K33 default)",
            font=MONO_FONT, size=16, bold=True, color=TEXT_DARK)

add_textbox(s, Inches(8.7), Inches(4.4), Inches(4.0), Inches(0.4),
            "PORT (Windows)", font=BODY_FONT, size=11, bold=True,
            color=TEAL)
add_textbox(s, Inches(8.7), Inches(4.8), Inches(4.0), Inches(0.5),
            "COM4  ·  9600 baud",
            font=MONO_FONT, size=16, bold=True, color=TEXT_DARK)

add_textbox(s, Inches(8.7), Inches(5.6), Inches(4.0), Inches(0.5),
            "Set SERIAL_PORT in server/app.py to match your machine.",
            font=BODY_FONT, size=11, italic=True, color=TEXT_MUTED)

add_footer(s, 4, TOTAL)


# ---------- Slide 5: Arduino firmware ----------
s = prs.slides.add_slide(blank_layout)
add_bg(s, WHITE)
add_eyebrow(s, "Firmware")
add_title(s, "Arduino reads sensors and emits JSON")

steps = [
    ("Read sensors",
     "DHT11 → temperature + humidity. TEMT6000 → analog light (0–1023)."),
    ("Update LED matrix",
     "Pick one of four icons based on the current humidity and light."),
    ("Send JSON over Serial",
     "Single-line JSON every 3 s. Flask parses it on the laptop side."),
]
for i, (title, body) in enumerate(steps):
    y = 1.85 + i * 1.3
    add_numbered_circle(s, 1.05, y + 0.45, 0.7, i + 1, color=TEAL)
    add_textbox(s, Inches(1.65), Inches(y), Inches(5.4), Inches(0.5),
                title, font=HEADER_FONT, size=19, bold=True,
                color=TEXT_DARK)
    add_textbox(s, Inches(1.65), Inches(y + 0.5), Inches(5.4),
                Inches(0.8), body, font=BODY_FONT, size=14,
                color=TEXT_MUTED)

# Code block on right
code_block(s, 7.2, 1.8, 5.6, 4.6, [
    '// loop(), every 3 seconds',
    'float h = dht.readHumidity();',
    'float t = dht.readTemperature();',
    'int   l = analogRead(A0);',
    '',
    'updateMatrixIcon(h, l);',
    '',
    'Serial.print("{\\"device_id\\":\\"Room_A01\\"");',
    'Serial.print(",\\"temperature\\":"); Serial.print(t, 1);',
    'Serial.print(",\\"humidity\\":");    Serial.print(h, 1);',
    'Serial.print(",\\"light\\":");       Serial.print(l);',
    'Serial.println("}");',
])

add_footer(s, 5, TOTAL)


# ---------- Slide 6: Flask server ----------
s = prs.slides.add_slide(blank_layout)
add_bg(s, WHITE)
add_eyebrow(s, "Server")
add_title(s, "Flask is the bridge between USB and the LAN")

# Left: three rows (icon + label + body)
bullets = [
    ("Listens on USB",
     "PySerial reads the COM port at 9600 baud and parses each JSON line."),
    ("Caches the latest reading",
     "Last good reading is held in memory; bad lines do not crash the server."),
    ("Exposes endpoints over LAN",
     "Binds to 0.0.0.0:5000 so any device on the same Wi-Fi can connect."),
]
for i, (title, body) in enumerate(bullets):
    y = 1.85 + i * 1.4
    add_circle(s, Inches(0.7), Inches(y + 0.2), Inches(0.4),
               color=SEAFOAM)
    add_textbox(s, Inches(1.3), Inches(y), Inches(5.6), Inches(0.5),
                title, font=HEADER_FONT, size=20, bold=True,
                color=TEXT_DARK)
    add_textbox(s, Inches(1.3), Inches(y + 0.5), Inches(5.6),
                Inches(0.9), body, font=BODY_FONT, size=14,
                color=TEXT_MUTED)

# Right: routes table
add_card(s, Inches(7.3), Inches(1.8), Inches(5.5), Inches(4.5),
         fill=SOFT_BG)
add_textbox(s, Inches(7.55), Inches(2.0), Inches(5.0), Inches(0.4),
            "ROUTES", font=BODY_FONT, size=11, bold=True, color=TEAL)

routes = [
    ("/",        "Main dashboard page"),
    ("/data",    "Live JSON, polled every 3 s"),
    ("/qr",      "Dynamic QR page"),
    ("/qr.png",  "Fresh QR image (in memory)"),
]
ry = 2.55
for path, desc in routes:
    add_textbox(s, Inches(7.55), Inches(ry), Inches(1.7), Inches(0.4),
                path, font=MONO_FONT, size=14, bold=True, color=TEAL)
    add_textbox(s, Inches(9.3), Inches(ry), Inches(3.4), Inches(0.4),
                desc, font=BODY_FONT, size=13, color=TEXT_DARK)
    ry += 0.85

add_footer(s, 6, TOTAL)


# ---------- Slide 7: Magic QR Refresh ----------
s = prs.slides.add_slide(blank_layout)
add_bg(s, WHITE)
add_eyebrow(s, "Magic QR Refresh")
add_title(s, "A fresh QR every time the network changes")

# Problem / Solution split
add_textbox(s, Inches(0.6), Inches(1.85), Inches(6), Inches(0.4),
            "PROBLEM", font=BODY_FONT, size=11, bold=True,
            color=RGBColor(0xC0, 0x39, 0x2B))
add_textbox(s, Inches(0.6), Inches(2.2), Inches(6), Inches(0.6),
            "Laptop IPv4 changes with every hotspot.",
            font=HEADER_FONT, size=20, bold=True, color=TEXT_DARK)
add_textbox(s, Inches(0.6), Inches(2.85), Inches(6), Inches(1.0),
            "A static QR points to a stale address and the phone "
            "gets a “site can’t be reached” error.",
            font=BODY_FONT, size=14, color=TEXT_MUTED)

add_textbox(s, Inches(0.6), Inches(4.1), Inches(6), Inches(0.4),
            "SOLUTION", font=BODY_FONT, size=11, bold=True, color=TEAL)
add_textbox(s, Inches(0.6), Inches(4.45), Inches(6), Inches(0.6),
            "Auto-detect IPv4, regenerate QR per request.",
            font=HEADER_FONT, size=20, bold=True, color=TEXT_DARK)
add_textbox(s, Inches(0.6), Inches(5.1), Inches(6), Inches(1.0),
            "A UDP socket trick reveals the outgoing IP — no static "
            "config, no manual rebuild before each demo.",
            font=BODY_FONT, size=14, color=TEXT_MUTED)

# Code snippet on right top
code_block(s, 7.0, 1.85, 5.8, 2.4, [
    'def get_local_ipv4():',
    '    s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)',
    '    s.connect(("8.8.8.8", 80))',
    '    ip = s.getsockname()[0]',
    '    s.close()',
    '    return ip',
])

# QR preview on right bottom
qr_path = Path(__file__).parent / "docs" / "dashboard_qr.png"
if qr_path.exists():
    s.shapes.add_picture(str(qr_path), Inches(7.0), Inches(4.5),
                         height=Inches(2.0), width=Inches(2.0))
    add_textbox(s, Inches(9.2), Inches(4.6), Inches(3.6), Inches(0.4),
                "Generated in memory",
                font=BODY_FONT, size=11, bold=True, color=TEAL)
    add_textbox(s, Inches(9.2), Inches(4.95), Inches(3.6), Inches(0.5),
                "Live /qr.png", font=HEADER_FONT, size=18, bold=True,
                color=TEXT_DARK)
    add_textbox(s, Inches(9.2), Inches(5.45), Inches(3.6), Inches(1.0),
                "Encodes the current dashboard URL based on the "
                "detected IPv4. No file written to disk.",
                font=BODY_FONT, size=12, color=TEXT_MUTED)

add_footer(s, 7, TOTAL)


# ---------- Slide 8: mobile dashboard ----------
s = prs.slides.add_slide(blank_layout)
add_bg(s, WHITE)
add_eyebrow(s, "Mobile dashboard")
add_title(s, "Live numbers, polled every three seconds")

# Phone mock-up on left
phone_x, phone_y = 1.0, 1.8
phone_w, phone_h = 3.4, 5.2
phone = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(phone_x),
                           Inches(phone_y), Inches(phone_w),
                           Inches(phone_h))
fill_solid(phone, TEXT_DARK); no_line(phone)
try:
    phone.adjustments[0] = 0.08
except Exception:
    pass
# Screen
screen = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(phone_x + 0.15),
                            Inches(phone_y + 0.35),
                            Inches(phone_w - 0.3),
                            Inches(phone_h - 0.7))
fill_solid(screen, WHITE); no_line(screen)
try:
    screen.adjustments[0] = 0.04
except Exception:
    pass

# Header strip inside screen
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                           Inches(phone_x + 0.15),
                           Inches(phone_y + 0.35),
                           Inches(phone_w - 0.3),
                           Inches(0.5))
fill_solid(strip, TEAL); no_line(strip)
add_textbox(s, Inches(phone_x + 0.3),
            Inches(phone_y + 0.45),
            Inches(phone_w - 0.6), Inches(0.4),
            "Room A01 · live",
            font=BODY_FONT, size=12, bold=True, color=WHITE)

# Three stat rows
stats = [
    ("Temperature", "23.4", "°C", MINT),
    ("Humidity",    "47",   "%",  SEAFOAM),
    ("Light",       "612",  "/1023", TEAL),
]
y0 = phone_y + 1.05
for i, (label, value, unit, color) in enumerate(stats):
    yy = y0 + i * 1.0
    dot = add_circle(s, Inches(phone_x + 0.35),
                     Inches(yy + 0.25),
                     Inches(0.18), color=color)
    add_textbox(s, Inches(phone_x + 0.65), Inches(yy),
                Inches(phone_w - 0.8), Inches(0.35),
                label.upper(), font=BODY_FONT, size=10, bold=True,
                color=TEXT_MUTED)
    add_textbox(s, Inches(phone_x + 0.65), Inches(yy + 0.3),
                Inches(phone_w - 0.8), Inches(0.55),
                value, font=HEADER_FONT, size=28, bold=True,
                color=TEXT_DARK)
    add_textbox(s, Inches(phone_x + 1.95), Inches(yy + 0.55),
                Inches(1.0), Inches(0.35),
                unit, font=BODY_FONT, size=12, color=TEXT_MUTED)

# Status + updated at bottom
add_textbox(s, Inches(phone_x + 0.3),
            Inches(phone_y + phone_h - 1.05),
            Inches(phone_w - 0.6), Inches(0.3),
            "STATUS", font=BODY_FONT, size=10, bold=True,
            color=TEXT_MUTED)
add_textbox(s, Inches(phone_x + 0.3),
            Inches(phone_y + phone_h - 0.75),
            Inches(phone_w - 0.6), Inches(0.4),
            "Normal · updated 1 s ago",
            font=BODY_FONT, size=13, bold=True, color=TEAL)

# Right side explanation
explain = [
    ("/data returns JSON",
     "device_id, temperature, humidity, light, status, last_update."),
    ("Polled every 3 s",
     "Page uses fetch() and updates the DOM without reloading."),
    ("Mobile-first HTML",
     "Vanilla HTML, CSS, and JS — no frameworks, no build step."),
]
for i, (title, body) in enumerate(explain):
    y = 1.95 + i * 1.45
    add_circle(s, Inches(5.5), Inches(y + 0.18), Inches(0.35),
               color=MINT)
    add_textbox(s, Inches(6.05), Inches(y), Inches(7.0), Inches(0.5),
                title, font=HEADER_FONT, size=20, bold=True,
                color=TEXT_DARK)
    add_textbox(s, Inches(6.05), Inches(y + 0.5), Inches(7.0),
                Inches(0.9), body, font=BODY_FONT, size=14,
                color=TEXT_MUTED)

add_footer(s, 8, TOTAL)


# ---------- Slide 9: LED matrix icons ----------
s = prs.slides.add_slide(blank_layout)
add_bg(s, WHITE)
add_eyebrow(s, "Local feedback")
add_title(s, "Four icons on the 8×8 LED matrix")

# Render each icon as an 8x8 grid of small squares
ICONS = {
    "RAIN":  [(2,1),(3,1),(4,1),(1,2),(5,2),(0,3),(6,3),(0,4),(6,4),
              (1,5),(2,5),(3,5),(4,5),(5,5),(1,6),(3,6),(5,6),(1,7),(3,7),(5,7)],
    "MOON":  [(3,1),(2,2),(1,3),(2,3),(1,4),(2,4),(1,5),(2,5),(3,5),
              (2,6),(3,6),(4,6),(5,6)],
    "SUN":   [(3,1),(1,2),(3,2),(5,2),(2,3),(3,3),(4,3),(0,4),(2,4),(3,4),(4,4),(6,4),
              (2,5),(3,5),(4,5),(1,6),(3,6),(5,6),(3,7)],
    "SMILE": [(2,1),(3,1),(4,1),(1,2),(5,2),(0,3),(2,3),(4,3),(6,3),
              (0,4),(2,4),(4,4),(6,4),(0,5),(2,5),(4,5),(6,5),(1,6),(5,6),
              (2,7),(3,7),(4,7)],
}
LABELS = [
    ("RAIN",  "humidity > 70 %",  SEAFOAM),
    ("MOON",  "light < 150",      TEAL),
    ("SUN",   "light > 850",      MINT),
    ("SMILE", "all normal",       SEAFOAM),
]

cell = 0.22  # inches per LED
matrix_w = cell * 8
gap_x = 0.5
total_w = 4 * matrix_w + 3 * gap_x
left0 = (13.333 - total_w) / 2
top0 = 2.2

for i, (label, cond, color) in enumerate(LABELS):
    x0 = left0 + i * (matrix_w + gap_x)
    # Frame card around matrix
    frame = add_card(s, Inches(x0 - 0.18), Inches(top0 - 0.18),
                     Inches(matrix_w + 0.36), Inches(matrix_w + 0.36),
                     fill=DARK_BG, line=DARK_BG)
    # 8x8 grid
    for gy in range(8):
        for gx in range(8):
            on = (gx, gy) in ICONS[label]
            dot = s.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x0 + gx * cell + 0.02),
                Inches(top0 + gy * cell + 0.02),
                Inches(cell - 0.04),
                Inches(cell - 0.04))
            fill_solid(dot, color if on else RGBColor(0x22, 0x40, 0x46))
            no_line(dot)
    # Label below
    add_textbox(s, Inches(x0 - 0.3), Inches(top0 + matrix_w + 0.4),
                Inches(matrix_w + 0.6), Inches(0.4),
                label, font=HEADER_FONT, size=18, bold=True,
                color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(x0 - 0.3), Inches(top0 + matrix_w + 0.85),
                Inches(matrix_w + 0.6), Inches(0.4),
                cond, font=BODY_FONT, size=12, color=TEXT_MUTED,
                align=PP_ALIGN.CENTER)

# Priority caption
add_textbox(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.4),
            "Priority order: rain → moon → sun → smile.",
            font=BODY_FONT, size=14, italic=True, color=TEXT_MUTED,
            align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.4),
            "The matrix gives instant local feedback; the phone "
            "dashboard gives the full numerical picture.",
            font=BODY_FONT, size=14, italic=True, color=TEXT_MUTED,
            align=PP_ALIGN.CENTER)

add_footer(s, 9, TOTAL)


# ---------- Slide 10: demo flow ----------
s = prs.slides.add_slide(blank_layout)
add_bg(s, WHITE)
add_eyebrow(s, "Live demo")
add_title(s, "From plug-in to live dashboard in five steps")

steps = [
    ("Plug in the Arduino",
     "Sensors power on; matrix shows the start-up smile."),
    ("Start the Flask server",
     "python server/app.py — listens on COM4 and binds to :5000."),
    ("Open /qr on the laptop",
     "Magic QR Refresh detects the current IPv4 and renders a fresh code."),
    ("Scan the QR with the phone",
     "Browser opens the dashboard over the same Wi-Fi / hotspot."),
    ("Watch live readings update",
     "Numbers refresh every 3 s; matrix icon reflects the current state."),
]

n = len(steps)
top_step = 1.95
row_h = 0.95
for i, (title, body) in enumerate(steps):
    y = top_step + i * row_h
    add_numbered_circle(s, 1.05, y + 0.4, 0.7, i + 1,
                        color=TEAL if i % 2 == 0 else SEAFOAM)
    add_textbox(s, Inches(1.7), Inches(y), Inches(11), Inches(0.4),
                title, font=HEADER_FONT, size=19, bold=True,
                color=TEXT_DARK)
    add_textbox(s, Inches(1.7), Inches(y + 0.42), Inches(11),
                Inches(0.4), body, font=BODY_FONT, size=14,
                color=TEXT_MUTED)

add_footer(s, 10, TOTAL)


# ---------- Slide 11: conclusion ----------
s = prs.slides.add_slide(blank_layout)
add_bg(s, DARK_BG)
# Accent dot
add_circle(s, Inches(0.85), Inches(0.85), Inches(0.25), color=MINT)
add_textbox(s, Inches(1.25), Inches(0.78), Inches(8), Inches(0.5),
            "RESULTS & WRAP-UP",
            font=BODY_FONT, size=14, bold=True, color=MINT)

add_textbox(s, Inches(0.8), Inches(1.4), Inches(12), Inches(1.2),
            "A working end-to-end prototype.",
            font=HEADER_FONT, size=44, bold=True, color=WHITE)

# Two columns: done / out of scope
add_textbox(s, Inches(0.8), Inches(3.0), Inches(5.7), Inches(0.5),
            "WHAT WORKS TODAY",
            font=BODY_FONT, size=12, bold=True, color=MINT)
done_items = [
    "Live DHT11 + TEMT6000 readings",
    "JSON over USB serial → Flask",
    "/data API consumed by mobile dashboard",
    "Magic QR Refresh on every demo",
    "Four-icon LED matrix feedback",
]
ty = 3.55
for item in done_items:
    add_circle(s, Inches(0.85), Inches(ty + 0.12), Inches(0.18),
               color=MINT)
    add_textbox(s, Inches(1.15), Inches(ty), Inches(5.4), Inches(0.4),
                item, font=BODY_FONT, size=15,
                color=RGBColor(0xE0, 0xEF, 0xF0))
    ty += 0.5

add_textbox(s, Inches(7.0), Inches(3.0), Inches(5.7), Inches(0.5),
            "OUT OF SCOPE FOR NOW",
            font=BODY_FONT, size=12, bold=True, color=SEAFOAM)
out_items = [
    "Cloud or remote access",
    "Long-term data storage",
    "Authentication / user accounts",
    "Wi-Fi MCU redesign (ESP32, etc.)",
]
ty = 3.55
for item in out_items:
    ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.05),
                              Inches(ty + 0.12), Inches(0.18),
                              Inches(0.18))
    ring.fill.background()
    ring.line.color.rgb = SEAFOAM
    ring.line.width = Pt(1.5)
    add_textbox(s, Inches(7.35), Inches(ty), Inches(5.4), Inches(0.4),
                item, font=BODY_FONT, size=15,
                color=RGBColor(0xE0, 0xEF, 0xF0))
    ty += 0.5

add_textbox(s, Inches(0.8), Inches(6.75), Inches(12), Inches(0.4),
            "Thank you — questions welcome.",
            font=HEADER_FONT, size=18, italic=True,
            color=SEAFOAM)


# ---------- save ----------
out_path = Path(__file__).parent / "presentation.pptx"
prs.save(str(out_path))
print(f"Saved: {out_path}")
